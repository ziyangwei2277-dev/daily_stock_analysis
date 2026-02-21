# -*- coding: utf-8 -*-
import os
import io
import pandas as pd
from datetime import datetime
from pathlib import Path
from typing import Optional
from fastapi import FastAPI, Request, Form, File, UploadFile
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse

# 导入各类解析库
from docx import Document
import pdfplumber
from pptx import Presentation

from api.v1 import api_v1_router
from api.middlewares.error_handler import add_error_handlers
from api.v1.schemas.common import RootResponse, HealthResponse
from analyzer_service import analyze_stock 

def create_app(static_dir: Optional[Path] = None) -> FastAPI:
    app = FastAPI(title="Daily Stock Analysis API - Pro Loader", version="1.3.0")
    
    # 【终极 CORS 修复】允许所有来源
    app.add_middleware(
        CORSMiddleware,
        allow_origins=["*"],
        allow_credentials=True,
        allow_methods=["*"],
        allow_headers=["*"],
    )

    app.include_router(api_v1_router)
    add_error_handlers(app)

    @app.post("/api/v1/stock/analyze-with-file", tags=["Stock Analysis"])
    async def analyze_with_file(
        stock_code: str = Form(...),
        full_report: bool = Form(False),
        file: UploadFile = File(...)
    ):
        file_name = file.filename
        content_bytes = await file.read()
        file_text = f"--- 文件: {file_name} 内容提取 ---\n"
        
        try:
            ext = file_name.lower().split('.')[-1]
            # 1. 解析 Word (.docx)
            if ext == 'docx':
                doc = Document(io.BytesIO(content_bytes))
                file_text += "\n".join([para.text for para in doc.paragraphs])
            
            # 2. 解析 PDF (.pdf)
            elif ext == 'pdf':
                with pdfplumber.open(io.BytesIO(content_bytes)) as pdf:
                    file_text += "\n".join([page.extract_text() or "" for page in pdf.pages])
            
            # 3. 解析 PPT (.pptx)
            elif ext == 'pptx':
                prs = Presentation(io.BytesIO(content_bytes))
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            file_text += shape.text + "\n"
            
            # 4. 解析 Excel (.xlsx, .xls, .csv)
            elif ext in ['xlsx', 'xls', 'csv']:
                df = pd.read_excel(io.BytesIO(content_bytes)) if ext != 'csv' else pd.read_csv(io.BytesIO(content_bytes))
                file_text += df.to_markdown() # 转为 Markdown 格式方便 Gemini 3 阅读
            
            # 5. 解析纯文本 (.txt)
            else:
                file_text += content_bytes.decode("utf-8", errors="ignore")
                
        except Exception as e:
            file_text = f"格式解析失败 ({file_name}): {str(e)}"

        # 调用 Gemini 3 进行双源标注分析
        result = analyze_stock(
            stock_code=stock_code,
            full_report=full_report,
            reliable_file_content=file_text,
            reliable_file_name=file_name
        )

        return {
            "status": "success",
            "stock_code": stock_code,
            "analysis": getattr(result, 'analysis', "未生成结果") if result else "分析失败"
        }

    return app

app = create_app()
